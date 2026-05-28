"""
Generate a complete stablecoin lecture PPTX using python-pptx.
Reuses design patterns from the existing template.

Usage:
    uv run --with python-pptx python3 scripts/generate_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Design tokens ──────────────────────────────────────────────────
SLIDE_W = Inches(10.0)
SLIDE_H = Inches(5.6)

C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_GREEN      = RGBColor(0x00, 0xB8, 0x70)
C_LIGHT_GREEN = RGBColor(0x00, 0xE0, 0x8A)
C_DARK       = RGBColor(0x1A, 0x1A, 0x1A)
C_GRAY       = RGBColor(0x5F, 0x6B, 0x72)
C_GRAY_LIGHT = RGBColor(0x8C, 0x96, 0x9C)
C_BODY       = RGBColor(0x2A, 0x2A, 0x2A)
C_RED        = RGBColor(0xE2, 0x57, 0x4C)
C_BG_LIGHT   = RGBColor(0xF5, 0xF5, 0xF5)
C_TABLE_ALT  = RGBColor(0xF0, 0xFA, 0xF5)

FONT_PRIMARY = "Apple SD Gothic Neo"

page_number = 0  # global page counter


# ── Utility helpers ────────────────────────────────────────────────
def _set_font(run, size=12, color=C_BODY, bold=False, name=FONT_PRIMARY):
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = name


def _add_textbox(slide, left, top, width, height, text, size=12, color=C_BODY,
                 bold=False, alignment=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf.vertical_anchor = anchor
    except Exception:
        pass
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size, color, bold)
    return tf


def _add_paragraph(tf, text, size=12, color=C_BODY, bold=False,
                   space_before=Pt(4), space_after=Pt(2), alignment=PP_ALIGN.LEFT,
                   level=0):
    p = tf.add_paragraph()
    p.alignment = alignment
    p.level = level
    if space_before:
        p.space_before = space_before
    if space_after:
        p.space_after = space_after
    run = p.add_run()
    run.text = text
    _set_font(run, size, color, bold)
    return p


def _add_page_number(slide, num):
    _add_textbox(slide, Inches(9.2), Inches(5.15), Inches(0.6), Inches(0.3),
                 str(num), size=10, color=C_GRAY_LIGHT, alignment=PP_ALIGN.RIGHT)


def _fill_background(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_green_rect(slide, left, top, width, height, color=C_GREEN):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


# ── Slide builders ─────────────────────────────────────────────────

def add_title_slide(prs, title, subtitle):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _fill_background(slide, C_WHITE)

    # Green accent bar at top
    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))

    # Title
    _add_textbox(slide, Inches(0.8), Inches(1.5), Inches(8.4), Inches(1.5),
                 title, size=42, color=C_DARK, bold=True, alignment=PP_ALIGN.CENTER,
                 anchor=MSO_ANCHOR.MIDDLE)

    # Subtitle
    _add_textbox(slide, Inches(1.0), Inches(3.0), Inches(8.0), Inches(0.8),
                 subtitle, size=15, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    # Bottom bar
    _add_green_rect(slide, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06))
    _add_green_rect(slide, Inches(4.2), Inches(5.46), Inches(1.6), Inches(0.04), C_LIGHT_GREEN)

    _add_page_number(slide, page_number)
    return slide


def add_section_slide(prs, number, title, subtitle=""):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Green background
    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, SLIDE_H, C_GREEN)
    # Lighter accent strip
    _add_green_rect(slide, Inches(0), Inches(0), Inches(0.12), SLIDE_H, C_LIGHT_GREEN)

    # Large number
    _add_textbox(slide, Inches(0.8), Inches(0.8), Inches(3.0), Inches(2.5),
                 number, size=72, color=C_LIGHT_GREEN, bold=True)

    # Title
    _add_textbox(slide, Inches(0.8), Inches(2.8), Inches(8.4), Inches(1.2),
                 title, size=34, color=C_WHITE, bold=True)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Inches(0.8), Inches(3.9), Inches(8.4), Inches(0.7),
                     subtitle, size=14, color=RGBColor(0xCC, 0xFF, 0xE8))

    _add_page_number(slide, page_number)
    return slide


def add_content_slide(prs, breadcrumb, title, subtitle, body_items, highlight_first=False):
    """body_items: list of strings (bullet points)"""
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    # Green top accent
    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04))

    # Breadcrumb
    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(4.0), Inches(0.3),
                 breadcrumb, size=10, color=C_GRAY_LIGHT)

    # Title
    _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.6),
                 title, size=28, color=C_DARK, bold=True)

    # Subtitle
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.4),
                     subtitle, size=13, color=C_GRAY)

    # Body
    top_y = Inches(1.65) if subtitle else Inches(1.35)
    txBox = slide.shapes.add_textbox(Inches(0.6), top_y, Inches(8.8), Inches(3.6))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(body_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(5)
        p.space_after = Pt(3)
        run = p.add_run()
        if highlight_first and i == 0:
            run.text = item
            _set_font(run, 20, C_GREEN, bold=True)
        else:
            run.text = f"\u2022  {item}"
            _set_font(run, 12, C_BODY)

    _add_page_number(slide, page_number)
    return slide


def add_two_column_slide(prs, breadcrumb, title, subtitle, left_items, right_items,
                          left_title="", right_title=""):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04))

    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(4.0), Inches(0.3),
                 breadcrumb, size=10, color=C_GRAY_LIGHT)
    _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.6),
                 title, size=28, color=C_DARK, bold=True)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.4),
                     subtitle, size=13, color=C_GRAY)

    top_y = Inches(1.65) if subtitle else Inches(1.35)

    # Left column background
    _add_green_rect(slide, Inches(0.5), top_y, Inches(4.2), Inches(3.6),
                    RGBColor(0xF7, 0xF7, 0xF7))

    # Right column background
    _add_green_rect(slide, Inches(5.3), top_y, Inches(4.2), Inches(3.6),
                    RGBColor(0xF0, 0xFA, 0xF5))

    def _build_col(items, left, col_title, title_color):
        if col_title:
            _add_textbox(slide, left + Inches(0.2), top_y + Inches(0.1),
                         Inches(3.8), Inches(0.35),
                         col_title, size=13, color=title_color, bold=True)
        start_y = top_y + Inches(0.5) if col_title else top_y + Inches(0.15)
        txBox = slide.shapes.add_textbox(left + Inches(0.2), start_y,
                                          Inches(3.8), Inches(3.0))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            _set_font(run, 11, C_BODY)

    _build_col(left_items, Inches(0.5), left_title, C_GRAY)
    _build_col(right_items, Inches(5.3), right_title, C_GREEN)

    _add_page_number(slide, page_number)
    return slide


def add_table_slide(prs, breadcrumb, title, subtitle, headers, rows):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04))

    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(4.0), Inches(0.3),
                 breadcrumb, size=10, color=C_GRAY_LIGHT)
    _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.6),
                 title, size=28, color=C_DARK, bold=True)
    if subtitle:
        _add_textbox(slide, Inches(0.6), Inches(1.1), Inches(8.8), Inches(0.4),
                     subtitle, size=13, color=C_GRAY)

    top_y = Inches(1.65) if subtitle else Inches(1.35)
    num_rows = len(rows) + 1
    num_cols = len(headers)
    col_w = 8.8 / num_cols
    row_h = min(0.35, 3.5 / num_rows)

    table_shape = slide.shapes.add_table(
        num_rows, num_cols,
        Inches(0.6), top_y,
        Inches(8.8), Inches(row_h * num_rows)
    )
    table = table_shape.table

    # Header row
    for j, h in enumerate(headers):
        cell = table.cell(0, j)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                _set_font(run, 10, C_WHITE, bold=True)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_GREEN

    # Data rows
    for i, row in enumerate(rows):
        for j, val in enumerate(row):
            cell = table.cell(i + 1, j)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    _set_font(run, 9, C_BODY)
            if i % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = C_TABLE_ALT

    _add_page_number(slide, page_number)
    return slide


def add_highlight_slide(prs, breadcrumb, title, big_text, description, items=None):
    """Slide with a large highlighted number/text and supporting content."""
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04))

    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(4.0), Inches(0.3),
                 breadcrumb, size=10, color=C_GRAY_LIGHT)
    _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.6),
                 title, size=28, color=C_DARK, bold=True)

    # Big number
    _add_textbox(slide, Inches(0.6), Inches(1.4), Inches(8.8), Inches(1.0),
                 big_text, size=36, color=C_GREEN, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    # Description
    _add_textbox(slide, Inches(0.6), Inches(2.4), Inches(8.8), Inches(0.5),
                 description, size=13, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    if items:
        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(3.1), Inches(8.8), Inches(2.2))
        tf = txBox.text_frame
        tf.word_wrap = True
        for i, item in enumerate(items):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(4)
            p.space_after = Pt(2)
            run = p.add_run()
            run.text = f"\u2022  {item}"
            _set_font(run, 12, C_BODY)

    _add_page_number(slide, page_number)
    return slide


def add_quote_slide(prs, breadcrumb, title, quote, attribution=""):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.04))

    _add_textbox(slide, Inches(0.6), Inches(0.25), Inches(4.0), Inches(0.3),
                 breadcrumb, size=10, color=C_GRAY_LIGHT)
    _add_textbox(slide, Inches(0.6), Inches(0.55), Inches(8.8), Inches(0.6),
                 title, size=28, color=C_DARK, bold=True)

    # Green left bar for quote
    _add_green_rect(slide, Inches(1.0), Inches(1.8), Inches(0.06), Inches(2.0), C_GREEN)

    _add_textbox(slide, Inches(1.3), Inches(1.8), Inches(7.8), Inches(2.0),
                 f"\u201c{quote}\u201d", size=14, color=C_BODY)

    if attribution:
        _add_textbox(slide, Inches(1.3), Inches(4.0), Inches(7.8), Inches(0.4),
                     f"\u2014 {attribution}", size=11, color=C_GRAY)

    _add_page_number(slide, page_number)
    return slide


def add_closing_slide(prs, text, subtext=""):
    global page_number
    page_number += 1
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _fill_background(slide, C_WHITE)

    _add_green_rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06))
    _add_green_rect(slide, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06))

    _add_textbox(slide, Inches(1.0), Inches(1.8), Inches(8.0), Inches(1.5),
                 text, size=36, color=C_DARK, bold=True,
                 alignment=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    if subtext:
        _add_textbox(slide, Inches(1.0), Inches(3.3), Inches(8.0), Inches(0.6),
                     subtext, size=14, color=C_GRAY, alignment=PP_ALIGN.CENTER)

    _add_page_number(slide, page_number)
    return slide


# ── Build the presentation ─────────────────────────────────────────

def build_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # ================================================================
    # SLIDE 1: Title
    # ================================================================
    add_title_slide(
        prs,
        "스테이블코인",
        "스테이블코인 시대의 도래  |  과거, 현재, 미래 그리고 직접 체험"
    )

    # ================================================================
    # SLIDE 2: Speaker Intro
    # ================================================================
    add_two_column_slide(
        prs,
        breadcrumb="발표자 소개",
        title="발표자 소개",
        subtitle="Web2에서 Web3로 이직한 입장에서, 스테이블코인이야말로 Web2와 Web3가 만나는 접점이라고 느꼈습니다.",
        left_title="Web2",
        right_title="Web3",
        left_items=[
            "네이버",
            "쿠팡",
            "토스 간편결제",
            "토스증권",
            "삼쩜삼",
        ],
        right_items=[
            "현재 Web3 프로젝트",
        ],
    )

    # ================================================================
    # SLIDE 3: TOC
    # ================================================================
    add_content_slide(
        prs,
        breadcrumb="CONTENTS",
        title="목차",
        subtitle="",
        body_items=[
            "01  스테이블코인의 탄생과 시련 — 유래, 개념, 그리고 $40B이 사라진 날",
            "02  왜 스테이블코인이어야 하는가 — CBDC, 예금 토큰이 아닌 이유",
            "03  스테이블코인이 만드는 미래 — 기계가 돈을 쓰고, 국경이 사라지는 세상",
            "04  지금, 어디까지 왔는가 — $323B 시장, GENIUS Act, 그리고 한국",
            "05  직접 체험 — KSD 발행부터 x402 결제까지",
        ],
    )

    # ================================================================
    # SECTION 01: 과거
    # ================================================================
    # Slide 4: Section divider
    add_section_slide(prs, "01", "스테이블코인의 탄생과 시련",
                      "유래, 개념, 그리고 $40B이 사라진 날")

    # Slide 5: 스테이블코인이란?
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 정의",
        title="스테이블코인이란?",
        subtitle="블록체인 위에서 법정화폐(주로 달러)와 동일한 가치를 유지하도록 설계된 디지털 자산",
        body_items=[
            "1 USDC = 1 USD — 발행사가 달러 준비금을 보유하여 가치 보장",
            "블록체인의 속도·투명성 + 달러의 안정성을 결합",
            "24/7 전송 가능, 은행 계좌 불필요, 스마트 컨트랙트와 조합 가능",
            "2026년 5월 기준 전체 시총 $323B, 연간 거래량 $46T (Visa의 3배)",
        ],
        highlight_first=True,
    )

    # Slide 6: 왜 만들어졌는가
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 배경",
        title="왜 만들어졌는가?",
        subtitle="변동성의 저주 — 비트코인으로는 결제도, 거래도, 상거래도 불가능했다",
        body_items=[
            "2013~14년 BTC $1,000 돌파, 하지만 하루 30% 급등락",
            "커피를 사는 동안 가격이 바뀌고, 판매자가 가격을 책정할 수 없는 자산",
            '2012.1 J.R. Willett — "블록체인 위의 USDCoin" 아이디어 최초 제안',
            "필요했던 것: 블록체인의 속도 + 달러의 안정성",
        ],
    )

    # Slide 7: 타임라인
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 타임라인",
        title="스테이블코인 역사 타임라인",
        subtitle="2012년 아이디어에서 2026년 $323B 시장까지",
        body_items=[
            "2014.7  BitUSD — 최초 스테이블코인 (암호자산 담보형) → 실패",
            "2014.10 Tether(USDT) — 법정화폐 담보형, '은행에 $1 = USDT 1개'",
            "2015    이더리움 출시 — 스마트 컨트랙트 + ERC-20 → 프로그래머블 시대",
            "2017.12 MakerDAO/DAI — 최초 프로그래머블 스테이블코인",
            "2018.9  USDC 출시 — 투명성 강화, 월간 준비금 attestation",
            "2020.6  DeFi Summer — 스테이블코인이 DeFi의 기축 화폐로 자리잡음",
            "2022.5  Terra/UST 붕괴 ($40B 소멸)",
            "2025.7  GENIUS Act 서명 — 미국 최초 스테이블코인 법제화",
        ],
    )

    # Slide 8: 트릴레마 — 3가지 유형
    add_table_slide(
        prs,
        breadcrumb="01 과거 · 유형",
        title="스테이블코인의 3가지 유형",
        subtitle="각 유형은 서로 다른 트레이드오프를 가진다",
        headers=["유형", "원리", "대표", "장점", "위험"],
        rows=[
            ["법정화폐 담보형", "1코인당 $1 준비금", "USDC, USDT", "가장 안정적", "중앙화, 준비금 리스크"],
            ["암호자산 담보형", "ETH 등 과잉 담보", "DAI", "탈중앙화", "자본 비효율, 담보 변동"],
            ["알고리즘형", "알고리즘 공급량 조절", "UST (Terra)", "자본 효율적", "극도로 위험"],
        ],
    )

    # Slide 9: 법정화폐 담보형 상세
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 법정화폐 담보형",
        title="법정화폐 담보형",
        subtitle="가장 단순하고, 가장 널리 사용되는 모델",
        body_items=[
            "은행에 $1 입금 → USDT/USDC 1개 발행",
            "토큰 제출 → 은행에서 $1 인출 (환매)",
            "장점: 단순하고 직관적, 가장 안정적",
            "리스크: 발행사/은행 신뢰 필요, 준비금 투명성 이슈",
            "Tether: 분기 보고 (제한적) / Circle: 월간 attestation (투명)",
        ],
    )

    # Slide 10: 암호자산 담보형 상세
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 암호자산 담보형",
        title="암호자산 담보형",
        subtitle="코드만으로 발행·관리·청산이 이루어지는 달러",
        body_items=[
            "ETH를 스마트 컨트랙트에 150% 이상 과잉 담보 → DAI 발행",
            "담보 비율 최소선 이하 → 스마트 컨트랙트가 자동 청산",
            "장점: 탈중앙화, 사람의 개입 불필요",
            "리스크: 담보 가격 급락 시 연쇄 청산, 자본 비효율",
            "MakerDAO → Sky로 리브랜딩, Multi-Collateral DAI (2019)",
        ],
    )

    # Slide 11: 알고리즘형 상세
    add_content_slide(
        prs,
        breadcrumb="01 과거 · 알고리즘형",
        title="알고리즘형 스테이블코인",
        subtitle="담보 없이 알고리즘만으로 $1 페깅을 유지하려 한 시도",
        body_items=[
            "두 토큰(UST + LUNA) 간 차익거래 메커니즘으로 가격 유지",
            "UST 매도 → $1 어치 LUNA로 교환 → 차익거래",
            "Anchor Protocol: 연 20% 이자 — 지속 불가능한 보조금",
            "UST의 75%가 Anchor에 집중, 구조적 Bank Run 취약성",
            "결론: 구조적으로 실패가 예정된 모델이었다",
        ],
    )

    # Slide 12: Terra/UST 붕괴
    add_content_slide(
        prs,
        breadcrumb="01 과거 · Terra 붕괴",
        title="Terra/UST 붕괴 — 3일 만에 $40B 증발",
        subtitle="2022년 5월, 알고리즘형의 구조적 한계가 증명되다",
        body_items=[
            "5/7 (토) — UST $85M 대량 매도, $0.985 소폭 이탈",
            '5/8 (일) — LFG $1.5B BTC 투입, Do Kwon: "Deploying more capital—steady lads"',
            "5/9 (월) — Anchor $5B 인출 (전체의 35%), UST $0.35, Death Spiral 시작",
            "5/12-13 — LUNA $0.0001, UST $0.10 이하, 총 $40B+ 증발",
            "Death Spiral: UST 매도 → LUNA 발행 폭증 → LUNA 가격 폭락 → 차익거래 붕괴",
        ],
    )

    # Slide 13: Terra 피해 & 연쇄 파산
    add_content_slide(
        prs,
        breadcrumb="01 과거 · Terra 여파",
        title="개인 피해 & 연쇄 파산",
        subtitle="하나의 알고리즘 스테이블코인 붕괴가 업계 전체를 강타",
        body_items=[
            "Reddit r/TerraLuna에 자살예방 핫라인 고정 — 은퇴 자금, 결혼 자금 전액 소멸",
            "2022.6  Three Arrows Capital(3AC) $10B 파산 — Terra 노출",
            "2022.6  Celsius Network 거래 동결 → Chapter 11",
            "2022.7  Voyager Digital 파산 — 3AC에 $650M+ 대출 회수 불가",
            "2022.11 FTX/Alameda 붕괴 — 연쇄 파산으로 약해진 시장에 최후의 일격",
            'Do Kwon: 몬테네그로 체포 → 미국 유죄 인정 → 15년 실형 ("epic scale의 사기")',
        ],
    )

    # Slide 14: SVB/USDC 디페깅
    add_content_slide(
        prs,
        breadcrumb="01 과거 · SVB 디페깅",
        title="SVB 사태와 USDC 디페깅 (2023.3)",
        subtitle='"담보형도 안전하지 않다" — 준비금 $3.3B가 SVB에 예치되어 있었다',
        body_items=[
            "3/10 금 — SVB 파산, Circle 발표: 준비금 중 $3.3B가 SVB에",
            "3/11 새벽 — USDC $0.87 최저점, Curve 하루 거래량 $6.03B (역대 최대)",
            "3/11-12 — 차익거래자: $0.87에 매수 → $1.00 환매 = 15% 수익",
            "3/12 저녁 — 미 재무부+Fed+FDIC: SVB 예금 전액 보호",
            "3/13 월 — USDC $1.00 복귀",
            "사후: Circle 준비금 87%를 BlackRock 미국 단기 국채 펀드로 이전",
        ],
    )

    # Slide 15: Terra vs USDC 비교
    add_table_slide(
        prs,
        breadcrumb="01 과거 · 비교",
        title="Terra vs USDC — 같은 디페깅, 정반대의 결과",
        subtitle="실물 준비금의 유무가 운명을 갈랐다",
        headers=["", "Terra/UST", "SVB/USDC"],
        rows=[
            ["원인", "구조적 결함 (담보 없음)", "외부 사건 (은행 파산)"],
            ["차익거래", "메커니즘 자체 붕괴", "정상 작동 ($0.87→$1)"],
            ["환매", "불가능", "Circle 계속 보장"],
            ["결과", "영구 붕괴 ($0)", "3일 만에 복귀"],
            ["피해", "$40B+ 소멸", "일시적, 시총 $1.9B 감소"],
        ],
    )

    # ================================================================
    # SECTION 02: 왜 스테이블코인이어야 하는가
    # ================================================================
    # Slide 16: Section divider
    add_section_slide(prs, "02", "왜 스테이블코인이어야 하는가",
                      "CBDC, 예금 토큰이 아닌 이유")

    # Slide 17: vs 기존 결제
    add_content_slide(
        prs,
        breadcrumb="02 왜 스테이블코인 · 기존 결제",
        title="후보 1: 기존 결제 시스템",
        subtitle="카드, SWIFT, ACH — 사람이 은행 영업시간에 은행을 통해 돈을 보내는 구조",
        body_items=[
            "프로그래밍 불가 — Visa에 조건부 결제 코드를 짤 수 없다",
            "영업시간 제한 — SWIFT는 주말에 안 된다",
            "마이크로페이먼트 불가 — $0.001 결제 시 수수료가 결제보다 크다",
            "기계 접근 불가 — AI 에이전트는 카드를 만들 수 없다",
            "글로벌 호환성 부재 — 국가별 결제 시스템이 단절",
        ],
    )

    # Slide 18: vs CBDC
    add_content_slide(
        prs,
        breadcrumb="02 왜 스테이블코인 · CBDC",
        title="후보 2: CBDC (중앙은행 디지털 화폐)",
        subtitle="137개국 연구 중이나, 실제 성과는 저조",
        body_items=[
            "나이지리아 eNaira: 채택률 0.5%→6% (강제 정책 동원 후)",
            "중국 디지털 위안: 프로젝트 책임자 부패 혐의 해임",
            "허가형(Permissioned) — 중앙은행 통제, DeFi 조합 불가",
            "국가 단위 — 한국 CBDC는 한국에서만 사용 가능",
            "프라이버시 우려 — 정부의 실시간 거래 감시 가능",
            "CBDC = 디지털 지폐 / 스테이블코인 = 프로그래머블 달러",
        ],
    )

    # Slide 19: vs 예금 토큰
    add_content_slide(
        prs,
        breadcrumb="02 왜 스테이블코인 · 예금 토큰",
        title="후보 3: 예금 토큰 (Deposit Token)",
        subtitle="JPM Coin — 은행 예금을 블록체인 토큰으로 표현",
        body_items=[
            "접근성: JP모건 고객만 (KYC 필수) vs USDC 누구나 지갑만 있으면",
            "네트워크: 폐쇄형 Walled garden vs 오픈 네트워크, 모든 DeFi 조합",
            "크로스보더: JP모건 지점 국가만 vs 어디서든 즉시 전송",
            "AI 에이전트: 은행 계좌 필요 vs 코드로 직접 결제 가능",
            '예금 토큰 = "디지털 수표" / 스테이블코인 = "디지털 현금"',
        ],
    )

    # Slide 20: 4가지 속성 교차점
    add_content_slide(
        prs,
        breadcrumb="02 왜 스테이블코인 · 결론",
        title="법정화폐 담보형 스테이블코인이 우위인 이유",
        subtitle="4가지 속성의 교차점",
        body_items=[
            "① 안정성 — 진짜 달러/국채가 뒷받침 (알고리즘형과 차별화)",
            "② 개방성 — 누구나, 어디서나, 허가 없이 사용 (CBDC/예금 토큰과 차별화)",
            "③ 프로그래밍 — 스마트 컨트랙트와 조합 가능 (기존 결제와 차별화)",
            "④ 이미 작동 중 — 시총 $315B, 연간 $33T 거래 (이론이 아닌 현실)",
        ],
        highlight_first=False,
    )

    # Slide 21: 달러 패권 도구
    add_content_slide(
        prs,
        breadcrumb="02 왜 스테이블코인 · 달러 패권",
        title="보론: 스테이블코인은 달러 패권의 도구인가?",
        subtitle="Tether + Circle 미 국채 보유 $1,449B — 한국($120B)보다 12배",
        body_items=[
            "스테이블코인 발행 확대 → 준비금으로 미 국채 매입 → 국채 수요 증가 → 달러 가치 유지",
            '재무장관 베센트(2025.3): "스테이블코인 기술로 달러의 기축 통화 지위를 유지할 것"',
            "GENIUS Act: 준비금을 3개월 미만 미 국채로 제한 — 구조적 국채 수요 강제",
            "기축 통화 평균 수명 ~100년, 달러는 이미 105년째",
            "스테이블코인 = 미국의 디지털 달러화(Dollarization) 전략",
        ],
    )

    # ================================================================
    # SECTION 03: 미래
    # ================================================================
    # Slide 22: Section divider
    add_section_slide(prs, "03", "스테이블코인이 만드는 미래",
                      "기계가 돈을 쓰고, 국경이 사라지는 세상")

    # Slide 23: 국경 없는 송금
    add_table_slide(
        prs,
        breadcrumb="03 미래 · 송금",
        title="국경이 사라지는 송금",
        subtitle="$10,000 송금 비용 비교",
        headers=["방식", "수수료", "소요 시간", "절감률"],
        rows=[
            ["신용카드 (2.5%)", "$250", "1~3일", "-"],
            ["SWIFT (수수료+환전)", "$390", "3~5 영업일", "-"],
            ["스테이블코인 (~0.55%)", "$55", "수 초~3분", "85.9%"],
        ],
    )

    # Slide 24: 송금 — 이미 일어나는 변화
    add_content_slide(
        prs,
        breadcrumb="03 미래 · 송금",
        title="이미 일어나고 있는 변화",
        subtitle="14억 명의 unbanked에게 스마트폰만으로 달러를 제공",
        body_items=[
            "2025년 B2B 스테이블코인 결제: $226B (전년 대비 733% 성장)",
            "필리핀 OFW 송금: 수수료 6% → 1% (연 $38.3B 규모)",
            "케냐 프리랜서: 수수료 29% → 2% (Mercy Corps 파일럿)",
            "전 세계 14억 명이 은행 계좌 없이도 달러를 보유·전송 가능",
            '"달러 은행 계좌 없이도 달러를 보유할 수 있다" — 금융 포용의 핵심',
        ],
    )

    # Slide 25: AI 에이전트 경제
    add_highlight_slide(
        prs,
        breadcrumb="03 미래 · AI 에이전트",
        title="기계가 기계에게 돈을 내는 세상",
        big_text="1.76억 건  |  $73M  |  USDC 98.6%",
        description="2025.5~2026.4 AI 에이전트 트랜잭션 (평균 $0.31~$0.48)",
        items=[
            "AI 에이전트는 카드도, 계좌도 없다 — 코드로 직접 결제 가능한 화폐가 필요",
            "x402 프로토콜: HTTP 402 Payment Required + USDC 자동 결제",
            "x402 Foundation: Google, Visa, Stripe, Mastercard, AWS, Microsoft 참여",
            "Gartner: 2030년 기계 고객이 전체 구매의 20%, $30T 규모",
        ],
    )

    # Slide 26: x402 프로토콜 상세
    add_content_slide(
        prs,
        breadcrumb="03 미래 · x402",
        title="x402 프로토콜 — HTTP 네이티브 결제",
        subtitle="1997년부터 존재했지만 아무도 쓰지 않던 HTTP 402를 스테이블코인이 살렸다",
        body_items=[
            "AI 에이전트: GET /api/data",
            '서버: 402 Payment Required {"price": "$0.001", "token": "USDC"}',
            "AI 에이전트: 자동으로 USDC $0.001 결제",
            "서버: 200 OK + 데이터 제공",
            "계정 없음, 구독 없음, API 키 없음 — 결제만으로 접근",
            "x402 처리: 1.69억 건, 48만 에이전트 (2026.4 기준)",
        ],
    )

    # Slide 27: 스트리밍 급여 & Pay-Per-Use
    add_two_column_slide(
        prs,
        breadcrumb="03 미래 · 스트리밍/Pay-Per-Use",
        title="스트리밍 급여 & Pay-Per-Use",
        subtitle="일하는 매 순간 돈이 들어오고, 쓴 만큼만 정확히 지불하는 세상",
        left_title="스트리밍 급여",
        right_title="Pay-Per-Use",
        left_items=[
            "연봉 $60,000 = 초당 $0.0019",
            "09:00 출근 → 급여 스트림 시작",
            "12:00 오전 근무분 $95.04 수령",
            "Superfluid, Sablier, Zebec 운영 중",
            "Deel: 2026.5 스테이블코인 급여 정식 출시",
        ],
        right_items=[
            "넷플릭스 $15.49/월 → 영화 1편 $0.99",
            "NYT $120/년 → 기사 1개 $0.10",
            "Spotify $10.99/월 → 곡당 $0.005",
            "Cloudflare NetDollar: AI 크롤러 실시간 소액 결제",
            '"구독 피로" 없이, 쓴 만큼만 지불',
        ],
    )

    # Slide 28: RWA 토큰화
    add_content_slide(
        prs,
        breadcrumb="03 미래 · RWA",
        title="토큰화된 세상의 기축 화폐",
        subtitle="실물 자산 토큰화(RWA) 2026.5 기준 $32B+ (전년 대비 200% 성장)",
        body_items=[
            "토큰화된 미국 국채, 부동산, 주식, 채권, 미술품",
            "거래 시 결제 수단 = 스테이블코인",
            "Mastercard: 2026.3 스테이블코인 인프라 기업 BVNK 인수",
            "맨해튼 빌딩의 0.001% 지분을 $50에 구매, 매달 임대 수익 자동 입금",
            "최소 투자금 $1 — 부동산 투자에 수억원이 필요한 시대의 종말",
        ],
    )

    # Slide 29: 감시와 통제의 그림자
    add_two_column_slide(
        prs,
        breadcrumb="03 미래 · 그림자",
        title="감시와 통제의 도구?",
        subtitle="프로그래머블 머니는 프로그래머블 감시이기도 하다",
        left_title="자유의 도구",
        right_title="통제의 도구",
        left_items=[
            "14억 명에게 은행 없이 달러를 준다",
            "독재 국가에서 자산을 보호한다",
            "국경을 넘어 즉시 가치를 전송한다",
        ],
        right_items=[
            "blacklist 함수로 자산 원격 동결",
            "모든 거래가 블록체인에 영구 기록",
            "정부 요청 시 발행사가 협조",
            "FinCEN + OFAC AML/제재 준수 의무",
        ],
    )

    # ================================================================
    # SECTION 04: 현재
    # ================================================================
    # Slide 30: Section divider
    add_section_slide(prs, "04", "지금, 어디까지 왔는가",
                      "$323B 시장, GENIUS Act, 그리고 한국")

    # Slide 31: 시장 현황
    add_highlight_slide(
        prs,
        breadcrumb="04 현재 · 시장",
        title="시장 현황 (2026년 5월)",
        big_text="$323B  |  Visa의 3배  |  ChatGPT Moment",
        description="전체 스테이블코인 시총 $323B, 2025년 온체인 거래량 $46T",
        items=[
            "USDT $189.6B (58.3%) / USDC $77.6B (2025년 73% 성장)",
            "달러 표시 비중 99%, Q1 2026 암호화폐 거래의 75%가 스테이블코인",
            "Citi Institute: 금융 시장의 'ChatGPT 모멘트' — 지수적 성장 변곡점",
            "2030년 시총 전망: J.P.Morgan $500B / Goldman $1T / Citi $1.6T",
        ],
    )

    # Slide 32: Tether 수익 모델
    add_content_slide(
        prs,
        breadcrumb="04 현재 · Tether",
        title="Tether의 수익 모델 — 직원 100명, 순이익 $13B",
        subtitle="스테이블코인 발행사는 예금 비용 0%로 은행보다 2배 마진",
        body_items=[
            "① 투자 수익: 준비금(미 국채, BTC, 금) 이자 — 미 국채 $1,200B (세계 19위)",
            "② 수수료: 법정화폐→USDT 전환 0.1% (최소 거래 $100,000)",
            "③ 대출 수익: USDT 담보 대출 (투명성 우려)",
            "은행: 예금 비용 10% → 마진 10% / 스테이블코인: 예금 비용 0% → 마진 20%",
            "한국 은행 수익의 90%가 이자마진 — 구조적 위협",
        ],
    )

    # Slide 33: Yield-bearing + Trump USD1
    add_content_slide(
        prs,
        breadcrumb="04 현재 · Yield-bearing",
        title="이자형 스테이블코인 & Trump USD1",
        subtitle="이자를 주는 달러에 대한 수요 폭발",
        body_items=[
            "USDe (Ethena): 파생상품 펀딩비 기반, 수익률 10~30%+, 높은 리스크",
            "USDS (Sky/MakerDAO): 미 국채+암호자산 담보, ~5%, 상대적 안전",
            "BUIDL (BlackRock): 미 국채, ~4.5%, 적격투자자 한정, 가장 안전",
            "Trump USD1: World Liberty Financial, 시총 ~$3B, 트럼프 가문 수익 75%",
            "대통령 가문이 스테이블코인 발행 — 시장의 전략적 중요성을 상징",
        ],
    )

    # Slide 34: USDT vs USDC
    add_table_slide(
        prs,
        breadcrumb="04 현재 · USDT vs USDC",
        title="USDT vs USDC — 두 거인의 경쟁",
        subtitle="규제가 강화될수록 USDC에 유리 — MiCA 시행 후 유럽 USDC 거래량 337% 급증",
        headers=["", "USDT (Tether)", "USDC (Circle)"],
        rows=[
            ["시총", "$189.6B (1위)", "$77.6B (2위)"],
            ["2025년 성장률", "36%", "73%"],
            ["투명성", "분기 보고 (제한적)", "월간 attestation (투명)"],
            ["규제 준수", "비규제", "MiCA 최초 준수, GENIUS Act 준비"],
            ["강점", "시장 점유율, 유동성", "투명성, 규제 준수, 기관 신뢰"],
        ],
    )

    # Slide 35: GENIUS Act
    add_content_slide(
        prs,
        breadcrumb="04 현재 · GENIUS Act",
        title="GENIUS Act — 미국 최초 스테이블코인 연방법",
        subtitle="2025.7.18 서명 — 스테이블코인이 법적 자산 클래스가 되다",
        body_items=[
            "정의: 지급결제용 스테이블코인 — 고정 가치 유지, 환매 의무",
            "준비금: 1:1 고품질 유동 자산(HQLA) — 현금, 예금, 3개월 미만 국채",
            "공시: 월간 준비금 attestation 공개 의무",
            "이자 지급 금지 — 보유자에게 이자/수익 지급 불가",
            "AML/제재: BSA 금융기관 분류, 자금세탁방지 프로그램 의무",
            "시행: 2027.1.18 또는 최종 규칙 공포 후 120일",
        ],
    )

    # Slide 36: 글로벌 규제 비교
    add_table_slide(
        prs,
        breadcrumb="04 현재 · 글로벌 규제",
        title="글로벌 규제 비교",
        subtitle="한국만 전용 프레임워크가 없다",
        headers=["항목", "미국 (GENIUS)", "EU (MiCA)", "일본 (PSA)", "한국"],
        rows=[
            ["법적 형태", "스테이블코인 전용법", "포괄적 암호자산법", "기존 금융법 개정", "전용법 없음 (8개 법안 계류)"],
            ["발행 주체", "OCC 승인", "은행/전자화폐기관", "은행/자금이동/신탁", "미정"],
            ["준비금", "HQLA, 3개월 미만 국채", "30% 분리 보관", "은행 요구불예금만", "미정"],
            ["이자 지급", "금지", "금지", "1:1 비율 유지", "미정"],
        ],
    )

    # Slide 37: 한국 현황
    add_content_slide(
        prs,
        breadcrumb="04 현재 · 한국",
        title="한국: 스테이블코인 규제의 공백",
        subtitle="8개 법안이 국회에 계류 중, 원화 스테이블코인 발행이 정부 공약에 포함",
        body_items=[
            "가상자산이용자보호법(2024.7) — 거래소 이용자 보호만, 발행·유통 규제 없음",
            "핵심 쟁점: 누가 발행? (은행 vs 비은행) / 누가 감독? (금융위 vs 한은)",
            "이재명 정부 공약: 디지털자산 기본법, STO 법제화, 원화 스테이블코인",
            "한국은행: 민간 발행 우려, CBDC 하이브리드 모델 선호 (Project Hangang)",
            "국내 시총 약 5,778억원 (글로벌의 ~0.1%), NHN KCP 원화 스테이블코인 상표 등록",
        ],
    )

    # Slide 38: Circle 사업 구조
    add_content_slide(
        prs,
        breadcrumb="04 현재 · Circle",
        title="Circle — 디지털 시대의 은행",
        subtitle="2025.6 NYSE 상장 (CRCL), 첫날 +168%, 기업 가치 $16B+",
        body_items=[
            "수익 구조: 사용자 $1 입금 → USDC 발행 → 국채 펀드 예치 → 이자 수취",
            "2025년 매출 $2.75B (+64%), 준비금 이자 비중 ~90%",
            "파트너(Coinbase 등)에 이자의 63% 지급, 나머지 37%가 Circle 매출",
            "은행 대비: 대출 없음 → 부실 채권 리스크 없음 (완전준비금 1:1)",
            "사용자 이자 없음 → USDC를 보유하는 이유 = 프로그래밍 가능성 + 글로벌 접근성",
        ],
    )

    # Slide 39: Core takeaway
    add_quote_slide(
        prs,
        breadcrumb="04 현재 · 핵심",
        title="결국 신뢰다",
        quote="2022년은 시련의 해였고 (Terra, FTX), 2023년은 회복의 해였고 (SVB 극복), "
              "2024년은 규제의 해였고 (MiCA), 2025년은 제도화의 해였다 (GENIUS Act, Circle IPO). "
              "2026년은 스테이블코인이 일상에 스며드는 해가 되고 있다.",
    )

    # Slide 40: 토론 질문
    add_content_slide(
        prs,
        breadcrumb="04 현재 · 토론",
        title="토론 질문",
        subtitle="",
        body_items=[
            "한국에 원화 스테이블코인이 필요한가? 통화 주권 vs 금융 혁신의 균형은?",
            "스테이블코인의 blacklist 기능은 필요한 규제인가, 자유의 침해인가?",
            "AI 에이전트가 자율적으로 결제하는 세상에서, 책임은 누구에게 있는가?",
            "은행이 스테이블코인 발행사가 되면, 기존 금융 시스템과 어떻게 달라지는가?",
        ],
    )

    # ================================================================
    # SECTION 05: 실습
    # ================================================================
    # Slide 41: Section divider
    add_section_slide(prs, "05", "직접 체험",
                      "KSD 발행부터 x402 결제까지")

    # Slide 42: 실습 개요
    add_table_slide(
        prs,
        breadcrumb="05 실습 · 개요",
        title="8단계 실습 흐름",
        subtitle="스테이블코인의 전체 생애주기를 직접 체험",
        headers=["Step", "실습 내용", "연결 강의", "시간"],
        rows=[
            ["1", "지갑 연결 + 에어드랍", "블록체인 지갑 개념", "5분"],
            ["2", "KSD 발행 (Mint)", "법정화폐 담보형 발행 구조", "5분"],
            ["3", "준비금 대시보드", "준비금의 중요성", "5분"],
            ["4", "KSD 소각 (Burn)", "환매와 차익거래", "3분"],
            ["5", "페깅 시뮬레이션 + SVB", "디페깅과 복원 원리", "7분"],
            ["6", "실제 USDC Mint 추적", "온체인 활동", "5분"],
            ["7", "CCTPv2 크로스체인 전송", "국경 없는 결제", "10분"],
            ["8", "x402 결제로 수료 뱃지", "AI 에이전트 경제", "5분"],
        ],
    )

    # Slide 43: Thank you
    add_closing_slide(
        prs,
        "Thank You",
        "스테이블코인의 과거, 현재, 미래를 함께 탐험해주셔서 감사합니다",
    )

    return prs


# ── Main ───────────────────────────────────────────────────────────

def main():
    out_dir = os.path.join(os.path.dirname(os.path.dirname(__file__)), "slides")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "stablecoin-lecture.pptx")

    prs = build_presentation()
    prs.save(out_path)
    print(f"Generated {page_number} slides → {out_path}")


if __name__ == "__main__":
    main()
